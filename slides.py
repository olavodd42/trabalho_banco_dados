# AVISO: Este código exige o pacote 'python-pptx', que deve ser executado em um ambiente local Python com suporte a bibliotecas externas.
# Execute este código no seu ambiente local com: pip install python-pptx

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
except ModuleNotFoundError:
    raise ImportError("A biblioteca 'python-pptx' não está instalada. Use 'pip install python-pptx' em seu ambiente local.")

# Criar apresentação
prs = Presentation()

# Função auxiliar para adicionar bullet points
def add_bullets(slide, title_text, bullet_points):
    slide.shapes.title.text = title_text
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

# Slide 1: Título
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Análise de Mandatos Presidenciais"
slide1.placeholders[1].text = "Banco de Dados e Visualizações Estatísticas"

# Slide 2: Descrição da Solução Adotada
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
add_bullets(slide2, "Descrição da Solução Adotada", [
    "Estrutura em MySQL com 4 tabelas principais",
    "Relacionamentos via chaves estrangeiras",
    "Dados importados de CSV e pré-processados em Python",
    "Métricas de duração calculadas em Python"
])

# Slide 3: Esquema do Banco de Dados (DDL)
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
slide3.shapes.title.text = "Esquema do Banco de Dados (DDL)"
ddl_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
tf3 = ddl_box.text_frame
for line in [
    "CREATE TABLE Presidente (id_presidente INT PRIMARY KEY, nome VARCHAR(100));",
    "CREATE TABLE Vice       (id_vice INT PRIMARY KEY, nome_vice VARCHAR(100));",
    "CREATE TABLE Partido    (id_partido INT PRIMARY KEY, nome_partido VARCHAR(100));",
    "CREATE TABLE Mandato    (id_mandato INT PRIMARY KEY, id_presidente INT, id_partido INT, id_vice INT, data_inicio DATE, data_fim DATE, duracao_dias INT, eleicao VARCHAR(50));"
]:
    p = tf3.add_paragraph()
    p.text = line
    p.level = 0
    p.font.size = Pt(14)

# Slide 4: Modelo ER
slide4 = prs.slides.add_slide(prs.slide_layouts[5])
slide4.shapes.title.text = "Modelo ER"
img_path = 'modelo_er.png'  # Caminho da imagem local
try:
    slide4.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))
except FileNotFoundError:
    text_box = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = "[Imagem modelo_er.png não encontrada]"

# Slide 5: Criação e Alimentação do Banco
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
add_bullets(slide5, "Criação e Alimentação do Banco", [
    "Execução de scripts DDL para criar as tabelas",
    "LOAD DATA INFILE para importar CSVs",
    "Verificação via SELECT COUNT(*) e LIMIT",
    "Configuração de local_infile no MySQL"
])

# Slide 6: Consultas SQL e Visualizações
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
add_bullets(slide6, "Consultas SQL e Visualizações", [
    "Histograma de distribuição de duração em anos",
    "Evolução da duração média por década (line plot)",
    "Box-plot de variabilidade por partido"
])

# Slide 7: Execução de Inserções (Exemplos)
slide7 = prs.slides.add_slide(prs.slide_layouts[5])
slide7.shapes.title.text = "Execução de Inserções"
ins_box = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
ins_tf = ins_box.text_frame
for cmd in [
    "LOAD DATA INFILE '/var/lib/mysql/presidentes.csv' INTO TABLE Presidente;",
    "LOAD DATA INFILE '/var/lib/mysql/partidos.csv' INTO TABLE Partido;",
    "LOAD DATA INFILE '/var/lib/mysql/vices.csv' INTO TABLE Vice;",
    "LOAD DATA INFILE '/var/lib/mysql/mandatos.csv' INTO TABLE Mandato;"
]:
    p = ins_tf.add_paragraph()
    p.text = cmd
    p.level = 0
    p.font.size = Pt(12)

# Salvar o arquivo
file_path = 'apresentacao_mandatos_final.pptx'
prs.save(file_path)
print(f"Arquivo salvo em: {file_path}")
